from __future__ import annotations

import io
import json
import shutil
import subprocess
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

import api.office_documents as office_documents
from api.office_documents import (
    CLAIMED_OFFICE_EXTENSIONS,
    is_claimed_office_path,
    preview_office_document,
    save_office_document,
)


ROOT = Path(__file__).resolve().parents[1]
WORKSPACE_JS = (ROOT / "static" / "workspace.js").read_text(encoding="utf-8")
NODE = shutil.which("node")


def _simple_docx_bytes(*paragraphs: str) -> bytes:
    document = DocxDocument()
    for paragraph in paragraphs or ("",):
        document.add_paragraph(paragraph)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _rich_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Lead paragraph")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "left"
    table.cell(0, 1).text = "right"
    document.add_paragraph("Tail paragraph")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _formatted_docx_bytes() -> bytes:
    document = DocxDocument()
    run = document.add_paragraph().add_run("Styled text")
    run.bold = True
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _paragraph_style_docx_bytes(style: str) -> bytes:
    document = DocxDocument()
    paragraph = document.add_paragraph("Styled paragraph")
    properties = paragraph._p.get_or_add_pPr()
    paragraph_style = OxmlElement("w:pStyle")
    paragraph_style.set(qn("w:val"), style)
    properties.append(paragraph_style)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _header_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Body text")
    document.sections[0].header.paragraphs[0].text = "Header text"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _custom_section_docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_paragraph("Body text")
    document.sections[0].left_margin = Inches(2)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _simple_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "alpha"
    sheet["B1"] = "beta"
    sheet["A2"] = "gamma"
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _simple_pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Office preview"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _multi_slide_pptx_bytes() -> bytes:
    presentation = Presentation()
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text = f"Slide {index + 1}"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_office_registry_claims_docx_xlsx_pptx():
    assert CLAIMED_OFFICE_EXTENSIONS == {".docx", ".xlsx", ".pptx"}
    assert is_claimed_office_path("report.docx")
    assert is_claimed_office_path("budget.xlsx")
    assert is_claimed_office_path("deck.pptx")

    docx_preview = preview_office_document("report.docx", _simple_docx_bytes("one", "two"))
    xlsx_preview = preview_office_document("budget.xlsx", _simple_xlsx_bytes())
    pptx_preview = preview_office_document("deck.pptx", _simple_pptx_bytes())

    assert docx_preview["preview_kind"] == "office"
    assert xlsx_preview["preview_kind"] == "office"
    assert pptx_preview["preview_kind"] == "office"
    assert docx_preview["office_format"] == "docx"
    assert xlsx_preview["office_format"] == "xlsx"
    assert pptx_preview["office_format"] == "pptx"
    assert docx_preview["render_mode"] == "code"
    assert xlsx_preview["render_mode"] == "code"
    assert pptx_preview["render_mode"] == "code"
    assert docx_preview["editable"] is True
    assert xlsx_preview["editable"] is False
    assert pptx_preview["editable"] is False


def test_docx_paragraph_projection_round_trips_simple_documents():
    original_bytes = _simple_docx_bytes("alpha", "beta")
    preview = preview_office_document("story.docx", original_bytes)

    assert preview["editable"] is True
    assert preview["content"] == "alpha\nbeta"

    saved_preview, saved_bytes = save_office_document("story.docx", original_bytes, "alpha\nbeta\ngamma")
    round_trip = DocxDocument(io.BytesIO(saved_bytes))

    assert saved_preview["editable"] is True
    assert saved_preview["content"] == "alpha\nbeta\ngamma"
    assert [paragraph.text for paragraph in round_trip.paragraphs] == ["alpha", "beta", "gamma"]


def test_docx_preview_preserves_interleaved_table_order():
    preview = preview_office_document("report.docx", _rich_docx_bytes())

    assert preview["content"] == "Lead paragraph\nTable 1\nleft\tright\nTail paragraph"


def test_docx_with_explicit_normal_style_stays_editable():
    original_bytes = _paragraph_style_docx_bytes("Normal")
    preview = preview_office_document("styled.docx", original_bytes)

    assert preview["editable"] is True
    assert preview["content"] == "Styled paragraph"

    saved_preview, saved_bytes = save_office_document("styled.docx", original_bytes, "Edited paragraph")
    round_trip = DocxDocument(io.BytesIO(saved_bytes))

    assert saved_preview["editable"] is True
    assert saved_preview["content"] == "Edited paragraph"
    assert [paragraph.text for paragraph in round_trip.paragraphs] == ["Edited paragraph"]


def test_docx_with_non_default_paragraph_style_stays_preview_only():
    preview = preview_office_document("heading.docx", _paragraph_style_docx_bytes("Heading1"))

    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")


def test_docx_preview_truncation_disables_editing(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_DOCX_PREVIEW_BLOCKS", 1)
    monkeypatch.setattr(
        office_documents,
        "_docx_editability",
        lambda _document: pytest.fail("truncated docx preview should not run full editability scan"),
    )

    preview = preview_office_document("story.docx", _simple_docx_bytes("alpha", "beta"))

    assert preview["truncated"] is True
    assert preview["editable"] is False
    assert preview["edit_blocked_reason"] == "docx preview exceeds safe limits"
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_docx_preview_char_budget_disables_editing(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_OFFICE_PREVIEW_CHARS", 5)
    monkeypatch.setattr(
        office_documents,
        "_docx_editability",
        lambda _document: pytest.fail("char-budget truncation should not run full editability scan"),
    )

    preview = preview_office_document("story.docx", _simple_docx_bytes("alphabet", "beta"))

    assert preview["truncated"] is True
    assert preview["editable"] is False
    assert preview["edit_blocked_reason"] == "docx preview exceeds safe limits"
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def _office_state_block() -> str:
    marker = "if(data.preview_kind==='office'){"
    start = WORKSPACE_JS.find(marker)
    assert start >= 0, "office preview state block not found in static/workspace.js"
    depth = 0
    for idx in range(start, len(WORKSPACE_JS)):
        char = WORKSPACE_JS[idx]
        if char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                return WORKSPACE_JS[start:idx + 1]
    raise AssertionError("could not find balanced office preview state block")


def _run_office_state_block(payload: dict) -> dict:
    js = (
        "const data = " + json.dumps(payload) + ";\n"
        + "const path = 'report';\n"
        + "let _previewRawContent = null;\n"
        + "let _previewRawContentPath = null;\n"
        + "let _previewServerEditable = null;\n"
        + "let _previewPreviewKind = '';\n"
        + "let _previewOfficeFormat = '';\n"
        + "let _previewSaveRoute = '/api/file/save';\n"
        + _office_state_block()
        + "\nconsole.log(JSON.stringify({_previewRawContent,_previewRawContentPath,_previewServerEditable,_previewPreviewKind,_previewOfficeFormat,_previewSaveRoute}));\n"
    )
    result = subprocess.run([NODE, "-e", js], check=True, capture_output=True, text=True, timeout=30)
    return json.loads(result.stdout.strip())


def test_xlsx_stays_preview_only():
    path = "budget.xlsx"
    raw = _simple_xlsx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_xlsx_preview_is_bounded(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_XLSX_PREVIEW_ROWS_PER_SHEET", 1)

    preview = preview_office_document("budget.xlsx", _simple_xlsx_bytes())

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_xlsx_preview_stops_when_char_budget_is_exhausted(monkeypatch):
    class FakeSheet:
        title = "Summary"
        max_row = 9_999
        max_column = 9_999

        def iter_rows(self, *, values_only, max_row, max_col):
            assert values_only is True
            assert max_row == office_documents.MAX_XLSX_PREVIEW_ROWS_PER_SHEET
            assert max_col == office_documents.MAX_XLSX_PREVIEW_CELLS_PER_SHEET
            yield ("x" * 200, "y" * 200)
            pytest.fail("xlsx preview should stop after the preview budget is exhausted")

    class FakeWorkbook:
        def __init__(self):
            self.worksheets = [FakeSheet()]

        def close(self):
            return None

    monkeypatch.setattr(office_documents, "MAX_OFFICE_PREVIEW_CHARS", 32)
    monkeypatch.setattr(office_documents, "_load_workbook_reader", lambda: lambda *_args, **_kwargs: FakeWorkbook())

    preview = preview_office_document("budget.xlsx", b"placeholder")

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_pptx_stays_preview_only():
    path = "deck.pptx"
    raw = _simple_pptx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_pptx_preview_is_bounded(monkeypatch):
    monkeypatch.setattr(office_documents, "MAX_PPTX_PREVIEW_SLIDES", 1)

    preview = preview_office_document("deck.pptx", _multi_slide_pptx_bytes())

    assert preview["truncated"] is True
    assert office_documents.OFFICE_PREVIEW_TRUNCATED_NOTICE in preview["content"]


def test_rich_docx_stays_preview_only():
    path = "rich.docx"
    raw = _rich_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_formatted_docx_stays_preview_only():
    path = "formatted.docx"
    raw = _formatted_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")

def test_header_docx_stays_preview_only():
    path = "headed.docx"
    raw = _header_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


def test_custom_section_docx_stays_preview_only():
    path = "custom-section.docx"
    raw = _custom_section_docx_bytes()

    preview = preview_office_document(path, raw)

    assert preview["preview_kind"] == "office"
    assert preview["editable"] is False
    assert preview.get("edit_blocked_reason")

    with pytest.raises(ValueError):
        save_office_document(path, raw, "edited text")


@pytest.mark.skipif(NODE is None, reason="node not on PATH")
def test_workspace_office_state_block_routes_all_office_formats_through_office_save():
    xlsx_state = _run_office_state_block(
        {
            "preview_kind": "office",
            "office_format": "xlsx",
            "editable": False,
            "content": "sheet preview",
        }
    )
    docx_state = _run_office_state_block(
        {
            "preview_kind": "office",
            "office_format": "docx",
            "editable": True,
            "content": "doc preview",
        }
    )

    assert xlsx_state["_previewSaveRoute"] == "/api/file/office-save"
    assert xlsx_state["_previewServerEditable"] is False
    assert xlsx_state["_previewOfficeFormat"] == "xlsx"
    assert docx_state["_previewSaveRoute"] == "/api/file/office-save"
    assert docx_state["_previewServerEditable"] is True
